"""Native Office/PDF artifact creation and inspection."""

from __future__ import annotations

import html
import os
import tempfile
from pathlib import Path
from typing import Any


SUPPORTED_ARTIFACTS = {"docx", "xlsx", "pptx", "pdf"}


def _sections(content: dict[str, Any]) -> list[dict[str, Any]]:
    raw = content.get("sections") or []
    return [item for item in raw if isinstance(item, dict)]


def create_artifact(path: Path, kind: str, content: dict[str, Any]) -> dict[str, Any]:
    kind = kind.lower().lstrip(".")
    if kind not in SUPPORTED_ARTIFACTS:
        raise ValueError(f"Unsupported artifact format: {kind}")
    if path.suffix.lower() != f".{kind}":
        raise ValueError(f"Output path must end in .{kind}")
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, temporary = tempfile.mkstemp(prefix=f".{path.name}.", suffix=f".{kind}", dir=path.parent)
    os.close(fd)
    temp_path = Path(temporary)
    try:
        if kind == "docx":
            _create_docx(temp_path, content)
        elif kind == "xlsx":
            _create_xlsx(temp_path, content)
        elif kind == "pptx":
            _create_pptx(temp_path, content)
        else:
            _create_pdf(temp_path, content)
        # See fs_atomic.preserve_existing_metadata: os.replace() swaps
        # inodes, so regenerating an existing artifact would otherwise
        # silently drop its original mode/owner in favor of mkstemp's
        # 0600 + the running process's uid/gid.
        from .fs_atomic import preserve_existing_metadata
        preserve_existing_metadata(temp_path, path)
        os.replace(temp_path, path)
    finally:
        temp_path.unlink(missing_ok=True)
    return {
        "success": True, "artifact_type": kind, "path": str(path),
        "size_bytes": path.stat().st_size,
    }


def _create_docx(path: Path, content: dict[str, Any]) -> None:
    from docx import Document
    document = Document()
    title = str(content.get("title") or "")
    if title:
        document.add_heading(title, 0)
    for section in _sections(content):
        heading = str(section.get("heading") or "")
        if heading:
            document.add_heading(heading, level=min(max(int(section.get("level") or 1), 1), 9))
        body = section.get("content") or section.get("body") or ""
        paragraphs = body if isinstance(body, list) else str(body).split("\n\n")
        for paragraph in paragraphs:
            document.add_paragraph(str(paragraph))
    document.core_properties.title = title
    document.save(path)


def _safe_cell(value: Any, allow_formulas: bool) -> Any:
    if not allow_formulas and isinstance(value, str) and value.startswith(("=", "+", "-", "@")):
        return "'" + value
    return value


def _create_xlsx(path: Path, content: dict[str, Any]) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font
    workbook = Workbook()
    workbook.remove(workbook.active)
    allow_formulas = bool(content.get("allow_formulas", False))
    sheets = content.get("sheets") or [{"name": "Sheet1", "rows": content.get("rows") or []}]
    for index, spec in enumerate(sheets):
        if not isinstance(spec, dict):
            continue
        sheet = workbook.create_sheet(str(spec.get("name") or f"Sheet{index + 1}")[:31])
        rows = spec.get("rows") or []
        for row in rows:
            values = row if isinstance(row, list) else [row]
            sheet.append([_safe_cell(value, allow_formulas) for value in values])
        if rows and bool(spec.get("header", True)):
            for cell in sheet[1]:
                cell.font = Font(bold=True)
        sheet.freeze_panes = spec.get("freeze_panes")
        for column in sheet.columns:
            width = min(max((len(str(cell.value or "")) for cell in column), default=8) + 2, 60)
            sheet.column_dimensions[column[0].column_letter].width = width
    if not workbook.sheetnames:
        workbook.create_sheet("Sheet1")
    workbook.save(path)


def _create_pptx(path: Path, content: dict[str, Any]) -> None:
    from pptx import Presentation
    presentation = Presentation()
    title = str(content.get("title") or "")
    slides = content.get("slides") or []
    if title:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = str(content.get("subtitle") or "")
    for spec in slides:
        if not isinstance(spec, dict):
            continue
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(spec.get("title") or "")
        body = spec.get("body") or spec.get("bullets") or ""
        items = body if isinstance(body, list) else str(body).splitlines()
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, item in enumerate(items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(item)
    presentation.save(path)


def _create_pdf(path: Path, content: dict[str, Any]) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    styles = getSampleStyleSheet()
    story = []
    title = str(content.get("title") or "")
    if title:
        story.extend((Paragraph(html.escape(title), styles["Title"]), Spacer(1, 12)))
    for section in _sections(content):
        heading = str(section.get("heading") or "")
        if heading:
            story.append(Paragraph(html.escape(heading), styles["Heading2"]))
        body = section.get("content") or section.get("body") or ""
        paragraphs = body if isinstance(body, list) else str(body).split("\n\n")
        for paragraph in paragraphs:
            story.extend((Paragraph(html.escape(str(paragraph)).replace("\n", "<br/>"), styles["BodyText"]), Spacer(1, 8)))
    SimpleDocTemplate(str(path), pagesize=A4, title=title).build(story)


def inspect_artifact(path: Path, *, max_chars: int = 30_000) -> dict[str, Any]:
    kind = path.suffix.lower().lstrip(".")
    if kind == "docx":
        from docx import Document
        doc = Document(path)
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        details = {"paragraphs": len(doc.paragraphs)}
    elif kind == "xlsx":
        from openpyxl import load_workbook
        book = load_workbook(path, read_only=True, data_only=False)
        chunks = []
        details = {"sheets": book.sheetnames}
        for sheet in book.worksheets:
            chunks.append(f"## {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                chunks.append("\t".join("" if value is None else str(value) for value in row))
        text = "\n".join(chunks)
        book.close()
    elif kind == "pptx":
        from pptx import Presentation
        deck = Presentation(path)
        chunks = []
        for number, slide in enumerate(deck.slides, 1):
            chunks.append(f"## Slide {number}")
            chunks.extend(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text)
        text = "\n".join(chunks)
        details = {"slides": len(deck.slides)}
    elif kind == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        details = {"pages": len(reader.pages), "metadata": {str(k): str(v) for k, v in (reader.metadata or {}).items()}}
    else:
        raise ValueError("inspect_artifact supports .docx, .xlsx, .pptx, and .pdf")
    return {
        "success": True, "artifact_type": kind, "path": str(path), **details,
        "text": text[:max_chars], "truncated": len(text) > max_chars,
    }
